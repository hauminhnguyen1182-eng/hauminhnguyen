from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import os
import re

class DocReader:
    def read(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.docx':
            return self._read_word(file_path)
        elif ext == '.xlsx':
            return self._read_excel(file_path)
        elif ext == '.pptx':
            return self._read_pptx(file_path)
        return {"type": "unknown", "content": "", "metadata": {}, "structure": []}

    def read_template(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.docx':
            return self._read_word_template(file_path)
        elif ext == '.pptx':
            return self._read_pptx_template(file_path)
        return {"type": "unknown", "structure": [], "placeholders": []}

    def _read_word(self, path):
        doc = Document(path)
        content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        return {"type": "word", "content": content, "metadata": {"paragraphs": len(doc.paragraphs)}, "structure": []}

    def _read_word_template(self, path):
        doc = Document(path)
        structure = []
        placeholders = []

        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                style = para.style.name if para.style else "Normal"
                is_heading = "Heading" in style

                section = {
                    "index": i,
                    "text": para.text.strip(),
                    "style": style,
                    "is_heading": is_heading,
                    "level": int(re.search(r'\d', style).group()) if re.search(r'\d', style) else 0
                }
                structure.append(section)

                # Extract placeholders like [Tên công ty], [Địa chỉ], etc.
                found_placeholders = re.findall(r'\[([^\]]+)\]', para.text)
                placeholders.extend(found_placeholders)

        return {
            "type": "word",
            "structure": structure,
            "placeholders": list(set(placeholders)),
            "total_paragraphs": len(doc.paragraphs)
        }

    def _read_excel(self, path):
        wb = load_workbook(path)
        content = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            content += f"Sheet: {sheet}\n"
            for row in ws.iter_rows(values_only=True):
                content += " | ".join([str(c) if c else "" for c in row]) + "\n"
        return {"type": "excel", "content": content, "metadata": {"sheets": wb.sheetnames}, "structure": []}

    def _read_pptx(self, path):
        prs = Presentation(path)
        content = ""
        for i, slide in enumerate(prs.slides):
            content += f"Slide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content += shape.text + "\n"
        return {"type": "ppt", "content": content, "metadata": {"slides": len(prs.slides)}, "structure": []}

    def _read_pptx_template(self, path):
        prs = Presentation(path)
        structure = []
        placeholders = []

        for i, slide in enumerate(prs.slides):
            slide_data = {"index": i, "shapes": []}
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shape_data = {
                        "text": shape.text.strip(),
                        "shape_type": str(shape.shape_type)
                    }
                    slide_data["shapes"].append(shape_data)

                    found_placeholders = re.findall(r'\[([^\]]+)\]', shape.text)
                    placeholders.extend(found_placeholders)

            structure.append(slide_data)

        return {
            "type": "ppt",
            "structure": structure,
            "placeholders": list(set(placeholders)),
            "total_slides": len(prs.slides)
        }

if __name__ == "__main__":
    reader = DocReader()
    print("✅ DocReader ready")
