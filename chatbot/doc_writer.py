from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill

PPTX_COLORS = {
    "white": PptxRGB(255, 255, 255),
    "blue": PptxRGB(0, 120, 215),
    "pink": PptxRGB(255, 182, 193),
    "yellow": PptxRGB(255, 255, 153),
    "dark": PptxRGB(51, 51, 51),
}

class DocWriter:
    def create_word(self, content, output_path, title="Document"):
        doc = Document()
        style = doc.styles['Normal']
        style.font.name = 'Arial'
        style.font.size = Pt(11)

        heading = doc.add_heading(title, 0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        for line in content.split('\n'):
            if line.strip():
                doc.add_paragraph(line)

        doc.save(output_path)
        return output_path

    def create_pptx(self, content, output_path, title="Presentation"):
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PPTX_COLORS["blue"]

        txBox = slide.shapes.add_textbox(PptxInches(1), PptxInches(2.5), PptxInches(11), PptxInches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PptxPt(44)
        p.font.color.rgb = PPTX_COLORS["white"]
        p.font.bold = True

        sections = content.split('\n\n')
        for section in sections[:10]:
            if section.strip():
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.5), PptxInches(12), PptxInches(6))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = section.strip()[:500]
                p.font.size = PptxPt(24)
                p.font.color.rgb = PPTX_COLORS["dark"]

        prs.save(output_path)
        return output_path

    def create_excel(self, content, output_path):
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"

        header_fill = PatternFill(start_color="0078D7", end_color="0078D7", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)

        for i, line in enumerate(content.split('\n')):
            if line.strip():
                cells = line.split('|')
                for j, cell in enumerate(cells):
                    cell_obj = ws.cell(row=i+1, column=j+1, value=cell.strip())
                    if i == 0:
                        cell_obj.fill = header_fill
                        cell_obj.font = header_font

        wb.save(output_path)
        return output_path

if __name__ == "__main__":
    writer = DocWriter()
    print("✅ DocWriter ready")
